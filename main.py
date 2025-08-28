from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
import httpx, os, uuid

app = FastAPI()
OUTPUT_DIR = "generated_ppts"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Serve frontend
app.mount("/static", StaticFiles(directory="static"), name="static")
@app.get("/", response_class=HTMLResponse)
async def homepage():
    with open("static/index.html", "r", encoding="utf-8") as f:
        return f.read()

@app.post("/generate-ppt")
async def generate_ppt(
    text_input: str = Form(...),               # changed from text_file
    template_file: UploadFile = File(...),
    guidance: str = Form(""),
    api_key: str = Form(...)
):
    try:
        input_text = text_input  # directly use the text from textarea

        # Save uploaded template
        template_filename = f"{uuid.uuid4().hex}_{template_file.filename}"
        template_path = os.path.join(OUTPUT_DIR, template_filename)
        with open(template_path, "wb") as f:
            f.write(await template_file.read())

        # Call LLM
        url = "https://aipipe.org/openai/v1/chat/completions"
        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {"role": "system", "content": "Split input text into slide titles and bullet points."},
                {"role": "user", "content": f"Guidance: {guidance}\n\nText:\n{input_text}"}
            ]
        }

        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                url, headers={"Authorization": f"Bearer {api_key}"}, json=payload
            )

        if resp.status_code != 200:
            return JSONResponse({"error": f"LLM error: {resp.text}"}, status_code=500)

        content = resp.json()["choices"][0]["message"]["content"]
        slides_text = [s.strip() for s in content.split("\n\n") if s.strip()]

        # Generate PPT
        prs = Presentation(template_path)
        for s in slides_text:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            parts = s.split("\n", 1)
            slide.shapes.title.text = parts[0].strip()
            if len(parts) > 1:
                slide.placeholders[1].text = parts[1].strip()

        output_filename = f"{uuid.uuid4().hex}_generated.pptx"
        output_file = os.path.join(OUTPUT_DIR, output_filename)
        prs.save(output_file)

        return FileResponse(output_file, filename="generated.pptx")

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
